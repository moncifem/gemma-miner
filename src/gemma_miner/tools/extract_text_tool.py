"""Universal text extractor.

`extract_text` takes a path (or raw bytes) and returns plain text, dispatching
on file extension first and on magic bytes second. Every parser is optional —
if the supporting library is missing, the tool returns a clear, actionable
error message naming the extra to install (`pip install gemma-miner[parsers]`).

Formats supported out of the box:

  Office:     .pdf .docx .pptx .xlsx .odt .rtf
  Web:        .html .htm .xml .svg
  Data:       .json .ndjson .jsonl .yaml .yml .toml .csv .tsv
  Archives:   .zip .tar .tar.gz .tgz .gz   (recurses on members)
  Ebooks:     .epub
  Plain:      .txt .md .log .py .* (anything else decodable as utf-8/latin-1)

Output is concatenated text. For tabular data (csv/xlsx) the output is a
human-readable rendering, not the raw bytes. The artifact field carries the
structured form (dict / list of dicts) when applicable.
"""

from __future__ import annotations

import csv
import gzip
import io
import json
import re
import tarfile
import zipfile
from html.parser import HTMLParser
from pathlib import Path
from typing import TYPE_CHECKING, Any

from gemma_miner.tools.base import Tool, ToolResult

if TYPE_CHECKING:
    from gemma_miner.state import AgentState


# ── helpers ────────────────────────────────────────────────────────────────


class _HtmlText(HTMLParser):
    def __init__(self):
        super().__init__()
        self._buf: list[str] = []
        self._skip = 0
        self._skip_tags = {"script", "style", "noscript"}

    def handle_starttag(self, tag, attrs):
        if tag in self._skip_tags:
            self._skip += 1
        if tag in ("br", "p", "div", "li", "tr"):
            self._buf.append("\n")

    def handle_endtag(self, tag):
        if tag in self._skip_tags and self._skip > 0:
            self._skip -= 1
        if tag in ("p", "div", "li", "tr"):
            self._buf.append("\n")

    def handle_data(self, data):
        if self._skip == 0:
            self._buf.append(data)

    def text(self) -> str:
        out = "".join(self._buf)
        return re.sub(r"\n{3,}", "\n\n", out).strip()


def _missing(pkg: str, install_extra: str = "parsers") -> str:
    return (
        f"ERROR: required package '{pkg}' not installed. "
        f"Install it with: pip install gemma-miner[{install_extra}]   (or: pip install {pkg})"
    )


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


# ── per-format extractors ──────────────────────────────────────────────────


def _from_pdf(data: bytes) -> tuple[str, Any]:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        return _missing("pypdf"), None
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages):
        try:
            pages.append(page.extract_text() or "")
        except Exception as e:  # noqa: BLE001
            pages.append(f"[page {i}: extraction error: {e}]")
    return "\n\n".join(pages), {"n_pages": len(pages)}


def _from_docx(data: bytes) -> tuple[str, Any]:
    try:
        import docx  # type: ignore
    except ImportError:
        return _missing("python-docx"), None
    doc = docx.Document(io.BytesIO(data))
    paras = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            paras.append("\t".join(c.text for c in row.cells))
    return "\n".join(paras), {"n_paragraphs": len(doc.paragraphs)}


def _from_pptx(data: bytes) -> tuple[str, Any]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return _missing("python-pptx"), None
    prs = Presentation(io.BytesIO(data))
    chunks = []
    for i, slide in enumerate(prs.slides):
        chunks.append(f"--- slide {i + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks), {"n_slides": len(prs.slides)}


def _from_xlsx(data: bytes) -> tuple[str, Any]:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        return _missing("openpyxl"), None
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets: dict[str, list[list]] = {}
    chunks = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = [[c for c in row] for row in ws.iter_rows(values_only=True)]
        sheets[name] = rows
        chunks.append(f"--- sheet: {name} ({len(rows)} rows) ---")
        for r in rows[:200]:
            chunks.append("\t".join("" if c is None else str(c) for c in r))
        if len(rows) > 200:
            chunks.append(f"... [{len(rows) - 200} more rows]")
    return "\n".join(chunks), {"sheets": {n: len(r) for n, r in sheets.items()}}


def _from_odt(data: bytes) -> tuple[str, Any]:
    # odt = zip of xml; pull content.xml and strip tags
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            with z.open("content.xml") as f:
                xml = f.read().decode("utf-8", errors="replace")
    except (KeyError, zipfile.BadZipFile) as e:
        return f"ERROR: not a valid ODT: {e}", None
    text = re.sub(r"<[^>]+>", " ", xml)
    text = re.sub(r"\s+", " ", text).strip()
    return text, None


def _from_rtf(data: bytes) -> tuple[str, Any]:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
    except ImportError:
        return _missing("striprtf"), None
    return rtf_to_text(_decode_text(data)), None


def _from_epub(data: bytes) -> tuple[str, Any]:
    try:
        from ebooklib import ITEM_DOCUMENT, epub  # type: ignore
    except ImportError:
        return _missing("EbookLib"), None
    book = epub.read_epub(io.BytesIO(data))
    chunks = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        p = _HtmlText()
        p.feed(item.get_content().decode("utf-8", errors="replace"))
        chunks.append(p.text())
    return "\n\n".join(chunks), {"n_items": len(chunks)}


def _from_html(data: bytes) -> tuple[str, Any]:
    p = _HtmlText()
    p.feed(_decode_text(data))
    return p.text(), None


def _from_xml(data: bytes) -> tuple[str, Any]:
    text = re.sub(r"<\?xml[^?]*\?>", "", _decode_text(data))
    text = re.sub(r"<!--.*?-->", "", text, flags=re.DOTALL)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text, None


def _from_json(data: bytes) -> tuple[str, Any]:
    try:
        obj = json.loads(_decode_text(data))
    except json.JSONDecodeError as e:
        return f"ERROR: invalid JSON: {e}", None
    return json.dumps(obj, indent=2, ensure_ascii=False), obj


def _from_jsonl(data: bytes) -> tuple[str, Any]:
    rows = []
    for ln in _decode_text(data).splitlines():
        ln = ln.strip()
        if not ln:
            continue
        try:
            rows.append(json.loads(ln))
        except json.JSONDecodeError:
            continue
    return json.dumps(rows[:200], indent=2, ensure_ascii=False), rows


def _from_yaml(data: bytes) -> tuple[str, Any]:
    try:
        import yaml  # type: ignore
    except ImportError:
        return _missing("PyYAML"), None
    obj = yaml.safe_load(_decode_text(data))
    return json.dumps(obj, indent=2, ensure_ascii=False), obj


def _from_toml(data: bytes) -> tuple[str, Any]:
    try:
        import tomllib  # py311+
    except ImportError:  # noqa: BLE001
        try:
            import tomli as tomllib  # type: ignore
        except ImportError:
            return _missing("tomli"), None
    obj = tomllib.loads(_decode_text(data))
    return json.dumps(obj, indent=2, ensure_ascii=False), obj


def _from_csv(data: bytes, sep: str = ",") -> tuple[str, Any]:
    text = _decode_text(data)
    reader = csv.reader(io.StringIO(text), delimiter=sep)
    rows = list(reader)
    chunks = []
    for r in rows[:500]:
        chunks.append("\t".join(r))
    if len(rows) > 500:
        chunks.append(f"... [{len(rows) - 500} more rows]")
    return "\n".join(chunks), rows


def _from_zip(data: bytes) -> tuple[str, Any]:
    chunks: list[str] = []
    members: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            for info in z.infolist():
                if info.is_dir():
                    continue
                members.append(info.filename)
                with z.open(info) as f:
                    inner = f.read()
                text, _ = _extract_bytes(inner, info.filename)
                chunks.append(f"\n=== {info.filename} ===\n{text}")
    except zipfile.BadZipFile as e:
        return f"ERROR: bad zip: {e}", None
    return "".join(chunks), {"members": members}


def _from_tar(data: bytes, mode: str = "r:*") -> tuple[str, Any]:
    chunks: list[str] = []
    members: list[str] = []
    try:
        with tarfile.open(fileobj=io.BytesIO(data), mode=mode) as t:
            for m in t.getmembers():
                if not m.isfile():
                    continue
                members.append(m.name)
                f = t.extractfile(m)
                if f is None:
                    continue
                inner = f.read()
                text, _ = _extract_bytes(inner, m.name)
                chunks.append(f"\n=== {m.name} ===\n{text}")
    except tarfile.TarError as e:
        return f"ERROR: bad tar: {e}", None
    return "".join(chunks), {"members": members}


def _from_gzip(data: bytes, name_hint: str = "") -> tuple[str, Any]:
    try:
        inner = gzip.decompress(data)
    except OSError as e:
        return f"ERROR: bad gzip: {e}", None
    inner_name = name_hint[:-3] if name_hint.endswith(".gz") else name_hint
    return _extract_bytes(inner, inner_name)


def _from_plain(data: bytes) -> tuple[str, Any]:
    return _decode_text(data), None


# ── dispatch ───────────────────────────────────────────────────────────────


_EXT_MAP = {
    ".pdf": _from_pdf,
    ".docx": _from_docx,
    ".pptx": _from_pptx,
    ".xlsx": _from_xlsx,
    ".xls":  _from_xlsx,
    ".odt":  _from_odt,
    ".rtf":  _from_rtf,
    ".epub": _from_epub,
    ".html": _from_html,
    ".htm":  _from_html,
    ".xml":  _from_xml,
    ".svg":  _from_xml,
    ".json": _from_json,
    ".ndjson": _from_jsonl,
    ".jsonl":  _from_jsonl,
    ".yaml": _from_yaml,
    ".yml":  _from_yaml,
    ".toml": _from_toml,
    ".csv":  _from_csv,
    ".tsv":  lambda d: _from_csv(d, sep="\t"),
    ".zip":  _from_zip,
    ".tar":  _from_tar,
    ".tgz":  lambda d: _from_tar(d, mode="r:gz"),
}


def _sniff(data: bytes) -> str | None:
    """Best-effort magic-byte sniffing for files with no extension."""
    if data.startswith(b"%PDF"):
        return ".pdf"
    if data.startswith(b"PK\x03\x04"):
        # zip-based: could be docx/pptx/xlsx/odt/epub; treat as zip; caller
        # will recurse and find the actual content
        return ".zip"
    if data.startswith(b"\x1f\x8b"):
        return ".gz"
    head = data[:512].lstrip().lower()
    if head.startswith(b"<?xml") or head.startswith(b"<svg"):
        return ".xml"
    if head.startswith(b"<!doctype html") or head.startswith(b"<html"):
        return ".html"
    if head.startswith(b"{") or head.startswith(b"["):
        return ".json"
    return None


def _extract_bytes(data: bytes, name: str) -> tuple[str, Any]:
    lname = name.lower()
    # multi-part suffixes first
    if lname.endswith(".tar.gz") or lname.endswith(".tgz"):
        return _from_tar(data, mode="r:gz")
    if lname.endswith(".tar.bz2"):
        return _from_tar(data, mode="r:bz2")
    if lname.endswith(".gz"):
        return _from_gzip(data, lname)

    ext = Path(lname).suffix
    fn = _EXT_MAP.get(ext)
    if fn is None:
        sniffed = _sniff(data)
        if sniffed and sniffed in _EXT_MAP:
            fn = _EXT_MAP[sniffed]
    if fn is None:
        return _from_plain(data)
    return fn(data)


# ── tool wrapper ───────────────────────────────────────────────────────────


class ExtractTextTool(Tool):
    name = "extract_text"
    summary_fields = ("text_chars", "n_pages")
    description = (
        "Extract plain text from any file the agent has cached or written. "
        "Dispatches by extension (then magic bytes if unknown). Handles: "
        "pdf, docx, pptx, xlsx, odt, rtf, epub, html, xml, json/jsonl, yaml, "
        "toml, csv/tsv, zip/tar/gz archives (recurses on members), and plain "
        "text. Returns the extracted text plus a small metadata header. Use "
        "this on every PDF/XML/CSV you download via http_get before feeding "
        "the body into extract_structured or dataset_append. If a parser "
        "library is missing, the error message tells you which package to "
        "pip install.\n\n"
        "Storage convention: when you call extract_text on an attachment, "
        "ALSO write the result to disk next to the original binary with a "
        ".txt suffix and a numbered name like "
        "`items/item_0001/attachment_01.txt` (the binary itself should be "
        "`items/item_0001/attachment_01.<ext>`). This pairing makes the "
        "extracted corpus iterable with a single glob."
    )
    args_schema = {
        "path": {
            "type": "string",
            "description": "File path (absolute, or relative to workdir).",
        },
        "max_chars": {
            "type": "integer",
            "default": 16000,
            "description": "Truncate the returned text to this many chars.",
        },
    }

    def run(self, args: dict, state: "AgentState") -> ToolResult:
        path = args.get("path")
        if not path:
            return ToolResult(output="ERROR: 'path' required", error=True)
        p = Path(path)
        if not p.is_absolute():
            p = Path(state.workdir) / p
        if not p.exists():
            return ToolResult(output=f"ERROR: not found: {p}", error=True)
        if not p.is_file():
            return ToolResult(output=f"ERROR: not a file: {p}", error=True)

        try:
            data = p.read_bytes()
        except Exception as e:  # noqa: BLE001
            return ToolResult(output=f"ERROR reading {p}: {e}", error=True)

        text, meta = _extract_bytes(data, p.name)
        is_error = text.startswith("ERROR:") if isinstance(text, str) else False

        max_chars = int(args.get("max_chars") or 16000)
        full_len = len(text)
        if full_len > max_chars:
            text = text[:max_chars] + f"\n... [truncated, total {full_len} chars]"

        header = (
            f"path: {p}\n"
            f"bytes: {len(data)}\n"
            f"detected_ext: {p.suffix.lower() or _sniff(data) or '(plain)'}\n"
            f"meta: {meta}\n"
            f"text_chars: {full_len}\n"
            "--- text ---\n"
        )
        return ToolResult(output=header + text, artifact=meta, error=is_error)
